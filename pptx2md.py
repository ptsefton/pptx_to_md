#!/usr/bin/env uv
import argparse
import subprocess
import os
import html
from pathlib import Path
from pptx import Presentation
import html
import math
import shutil
import commonmark
from datetime import date

import pymupdf as fitz
   
def get_text(slide):
    text = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
            text += "\n"
    return(text)



def save_images_from_pdf(out_dir, pdf_filename, formatString):
    os.chdir(out_dir)

    try:
        doc = fitz.open(pdf_filename)
        for page_num in range(len(doc)):
            svg_filename = formatString % page_num
            page = doc[page_num]
            svg_text = page.get_svg_image()
            
            with open(svg_filename, 'w') as f:
                f.write(svg_text)
                
            print(f"Converted page {page_num + 1} to {svg_filename}")
        doc.close()
        return  # Success, exit function
    except Exception as e:
        print(f"PyMuPDF conversion failed: {e}")
        
      


def parse_preso(powerpoint_file, use_svg=False):
    slug = powerpoint_file.stem
    pptx = powerpoint_file.name
    pdf = powerpoint_file.with_suffix('.pdf').name
    prs = Presentation(powerpoint_file)
    title = prs.core_properties.title
    slides = []
    for slide in prs.slides:
        s = {"text": "", "notes": ""}
        s["text"] = get_text(slide)
        if slide.has_notes_slide:
            s["notes"] = get_text(slide.notes_slide)
        slides.append(s)

    # Fix: img_prefix needs to be defined or passed as parameter
    img_prefix = ""  # Add this line or pass as parameter
    
    md = f"""---
    title: {slug}
    date: {date.today()}
    slug: {slug}
    category:
    author:
---

<a href="{img_prefix}{pdf}">PDF version</a> | <a href="{img_prefix}{pptx}">Powerpoint Version</a>


    """
    
    digits = math.floor(math.log10(len(slides))) + 1
    count = -1
    
    # Choose file extension based on format
    extension = "svg" 
    formatString = f"Slide%0{digits}d.{extension}"

    for slide in slides:
        count += 1
        # TODO make this a functions
        
        image_path = formatString % count
        # compute alt text separately to avoid backslashes inside f-string expressions
        alt_text = html.escape(slide["text"], quote=True).replace("\n", " :: ")
        md += f"""

<section typeof='http://purl.org/ontology/bibo/Slide'>
<img src='{img_prefix}{image_path}' alt='{alt_text}' title='Slide: {count}' border='1'  width='85%'/> 


{slide["notes"]}


</section>

"""
    return md, formatString


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument('filename',  default=None, type=Path, help='Name of input Powerpoint file')
    parser.add_argument('-i', '--img-prefix', default = '', help="String to put in front of image paths")
    parser.add_argument('-d', '--dir', default=None, type=Path, help="Output directory (created automatically if it doesn't exist)")
    parser.add_argument('-t', '--topdf', default = False, action="store_true", help="Convert DOCX version to PDF using openoffice (soffice)")
    parser.add_argument('-s', '--soffice', default = "soffice", help="soffice command")
    parser.add_argument('--html', default = False, action="store_true", help="Generate an index.html file from the markdown")

    args = vars(parser.parse_args())

    img_prefix = args["img_prefix"]
    powerpoint_file = args["filename"].resolve()
    base_path = powerpoint_file.parent;
    pdf_file = powerpoint_file.with_suffix(".pdf")
    
    # Use custom directory if provided, otherwise use default (powerpoint_file.stem)
    if args["dir"]:
        out_dir = args["dir"].resolve()
    else:
        out_dir = base_path / powerpoint_file.stem
    
    md_path = out_dir / "index.md"


    if not powerpoint_file.is_file():
        raise Exception(f"{powerpoint_file} file not found")

    out_dir.mkdir(parents=True, exist_ok=True)

    shutil.copyfile(powerpoint_file, out_dir / powerpoint_file.name)
    shutil.copyfile(pdf_file, out_dir / pdf_file.name)

    md, formatString = parse_preso(powerpoint_file)
    
    with open(md_path, 'w', encoding='utf-8') as o:
        print("writing", md_path)
        o.write(md)

    # Generate HTML if requested
    if args["html"]:
        html_path = out_dir / "index.html"
        parser = commonmark.Parser()
        renderer = commonmark.HtmlRenderer()
        ast = parser.parse(md)
        html_content = renderer.render(ast)
        
        # Wrap in a basic HTML template
        full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{powerpoint_file.stem}</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 1200px; margin: 0 auto; padding: 20px; }}
        img {{ max-width: 100%; height: auto; }}
        section {{ margin: 20px 0; padding: 20px; border: 1px solid #ddd; }}
    </style>
</head>
<body>
{html_content}
</body>
</html>"""
        
        with open(html_path, 'w', encoding='utf-8') as o:
            print("writing", html_path)
            o.write(full_html)

    if args["topdf"]:
        os.chdir(out_dir)
        convert_command =  f'{args["soffice"]} --headless --convert-to pdf "{powerpoint_file.name}"'
        print(convert_command);
        subprocess.Popen(convert_command,
                        shell=True,
                        stdout=subprocess.PIPE).communicate()

    save_images_from_pdf(out_dir, pdf_file.name, formatString)

